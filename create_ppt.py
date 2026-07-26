from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
import glob

def get_latest_image(prefix):
    # Finds the latest generated image matching the prefix
    search_path = r"C:\Users\User\.gemini\antigravity\brain\44bc7520-08b1-4bf2-bcd8-62a9add7da07\\" + prefix + "*.png"
    files = glob.glob(search_path)
    if not files: return None
    return sorted(files)[-1]

def create_presentation():
    prs = Presentation()

    navy = RGBColor(11, 19, 43)
    
    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Bhoomi Rakshak"
    title.text_frame.paragraphs[0].font.color.rgb = navy
    title.text_frame.paragraphs[0].font.bold = True
    subtitle.text = "Automated Crop Insurance Claim Intake & Duplicate Detection System"
    img = get_latest_image("hero_title")
    if img:
        slide.shapes.add_picture(img, Inches(3.0), Inches(5.0), width=Inches(4.0))

    # 2. Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Problem: Manual Verification Flaws"
    tf = slide.placeholders[1].text_frame
    slide.placeholders[1].width = Inches(4.5)
    tf.text = "Current Challenges in the Block Office:"
    tf.add_paragraph().text = "• Claims are verified one at a time in isolation."
    tf.add_paragraph().text = "• Exact duplicates easily slip through."
    tf.add_paragraph().text = "• Over-claiming extent goes unnoticed."
    tf.add_paragraph().text = "• Near-duplicates evade simple checks."
    img = get_latest_image("digital_ledger")
    if img: slide.shapes.add_picture(img, Inches(5.0), Inches(2.0), width=Inches(4.5))

    # 3. Solution Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Solution Architecture & Tech Stack"
    tf = slide.placeholders[1].text_frame
    slide.placeholders[1].width = Inches(4.5)
    tf.text = "A simple, explainable, locally-hosted web application."
    tf.add_paragraph().text = "• Backend: Python & Flask for the REST API."
    tf.add_paragraph().text = "• Database: SQLite (Raw SQL)."
    tf.add_paragraph().text = "• ML: Scikit-Learn Logistic Regression."
    tf.add_paragraph().text = "• Frontend: Vanilla HTML/CSS/JS with a 'Ledger' aesthetic."
    img = get_latest_image("architecture_nodes")
    if img: slide.shapes.add_picture(img, Inches(5.0), Inches(2.0), width=Inches(4.5))

    # 4. Database Schema
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Normalized Database Schema"
    tf = slide.placeholders[1].text_frame
    slide.placeholders[1].width = Inches(4.5)
    tf.text = "Key Entities:"
    tf.add_paragraph().text = "• survey_numbers: The ground-truth village register."
    tf.add_paragraph().text = "• claims: Incoming claims to be verified."
    tf.add_paragraph().text = "• flags: Stores specific human-readable reasons."
    tf.add_paragraph().text = "• decisions: Enforces a strict UNIQUE constraint."
    img = get_latest_image("database_schema")
    if img: slide.shapes.add_picture(img, Inches(5.0), Inches(2.0), width=Inches(4.5))

    # 5. Rule-Based Checks
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Automated Rule-Based Checks"
    tf = slide.placeholders[1].text_frame
    slide.placeholders[1].width = Inches(4.5)
    tf.text = "Runs instantly at the counter upon submission:"
    tf.add_paragraph().text = "1. Exact Duplicate: Rejects if an active claim already exists."
    tf.add_paragraph().text = "2. Extent Exceeds: Flags if claimed acreage > ground-truth acreage."
    tf.add_paragraph().text = "3. Crop Mismatch: Flags if the claimed crop differs."
    img = get_latest_image("agricultural_field")
    if img: slide.shapes.add_picture(img, Inches(5.0), Inches(2.0), width=Inches(4.5))

    # 6. Near-Duplicate ML Model
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "ML: Near-Duplicate Detection"
    tf = slide.placeholders[1].text_frame
    slide.placeholders[1].width = Inches(4.5)
    tf.text = "Algorithm: Logistic Regression (class_weight='balanced')"
    tf.add_paragraph().text = "Engineered Features:"
    tf.add_paragraph().text = "• String Similarity: Jaro-Winkler distance on Name and Survey."
    tf.add_paragraph().text = "• Transposition Boolean: Explicitly checks for transposed digits."
    tf.add_paragraph().text = "• Extent Closeness: Normalized mathematical difference in acreage."
    img = get_latest_image("ml_network")
    if img: slide.shapes.add_picture(img, Inches(5.0), Inches(2.0), width=Inches(4.5))

    # 7. Model Performance
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "ML Model Performance"
    tf = slide.placeholders[1].text_frame
    slide.placeholders[1].width = Inches(4.5)
    tf.text = "Evaluated using rigorous Leave-One-Out Cross-Validation:"
    tf.add_paragraph().text = "• Recall: 100% (Caught 5 out of 5 planted near-duplicate claims)."
    tf.add_paragraph().text = "• False Positives: Only 42 out of 1,760 non-duplicate pairs scanned."
    tf.add_paragraph().text = "• Explainability: Shows the exact confidence score."
    img = get_latest_image("ml_performance")
    if img: slide.shapes.add_picture(img, Inches(5.0), Inches(2.0), width=Inches(4.5))

    # 8. Verification Queue
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Verification Queue"
    tf = slide.placeholders[1].text_frame
    slide.placeholders[1].width = Inches(4.5)
    tf.text = "Officer Dashboard Workflow:"
    tf.add_paragraph().text = "• Queue: Displays all flagged claims chronologically."
    tf.add_paragraph().text = "• Optimistic Locking: Prevents race conditions."
    tf.add_paragraph().text = "• Comparison: UI renders a side-by-side view."
    tf.add_paragraph().text = "• Decision: Requires a mandatory remark."
    img = get_latest_image("dashboard_ui")
    if img: slide.shapes.add_picture(img, Inches(5.0), Inches(2.0), width=Inches(4.5))

    # 9. Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Conclusion"
    tf = slide.placeholders[1].text_frame
    slide.placeholders[1].width = Inches(4.5)
    tf.text = "Bhoomi Rakshak achieves:"
    tf.add_paragraph().text = "• Data Integrity: Strict SQLite enforcement."
    tf.add_paragraph().text = "• Operational Efficiency: Live ground-truth lookups prevent errors."
    tf.add_paragraph().text = "• Fraud Prevention: Blends deterministic rules with statistical ML."
    img = get_latest_image("modern_office")
    if img: slide.shapes.add_picture(img, Inches(5.0), Inches(2.0), width=Inches(4.5))

    # Save
    out_name = "Bhoomi_Rakshak_Presentation_v3.pptx"
    prs.save(out_name)
    print(f"Presentation saved as {out_name}")

if __name__ == "__main__":
    create_presentation()
